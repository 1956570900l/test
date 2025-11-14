# backend/data_process/pdf_processor.py
# (V4 - 最终递归版，支持子文件夹)

import os
import fitz  # PyMuPDF
from PIL import Image
import io
from typing import List, Tuple, Dict
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from tqdm import tqdm 
from utils import config
import docx
import pptx

# --- 文本切分器 (全局) ---
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000, 
    chunk_overlap=100
)

# --- (不变) DOCX 文件处理器 ---
def process_docx(file_path: str, filename: str) -> List[Document]:
    print(f"  > (DOCX) 正在处理 {filename}")
    try:
        doc = docx.Document(file_path)
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text])
        if not full_text: return []
        
        chunks = text_splitter.split_text(full_text)
        doc_list = []
        for chunk in chunks:
            doc_list.append(Document(
                page_content=chunk,
                metadata={"doc_name": filename, "page": 1, "clause_id": ""}
            ))
        return doc_list
    except Exception as e:
        print(f"  > 错误: 处理 {filename} 失败: {e}")
        return []

# --- (不变) PPTX 文件处理器 ---
def process_pptx(file_path: str, filename: str) -> List[Document]:
    print(f"  > (PPTX) 正在处理 {filename}")
    try:
        prs = pptx.Presentation(file_path)
        full_text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = f"\n--- [幻灯片 {i+1}] ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"): slide_text += shape.text + "\n"
            full_text += slide_text
        
        if not full_text: return []

        chunks = text_splitter.split_text(full_text)
        doc_list = []
        for chunk in chunks:
            doc_list.append(Document(
                page_content=chunk,
                metadata={"doc_name": filename, "page": 1, "clause_id": ""}
            ))
        return doc_list
    except Exception as e:
        print(f"  > 错误: 处理 {filename} 失败: {e}")
        return []

# --- (不变) PDF 文件处理器 ---
def process_pdf(file_path: str, filename: str, output_image_dir: str) -> Tuple[List[Document], List[Dict]]:
    print(f"  > (PDF) 正在图文处理 {filename}")
    text_docs: List[Document] = []
    image_infos: List[Dict] = []
    
    try:
        doc = fitz.open(file_path)
        for page_num in tqdm(range(len(doc)), desc=f"    > 遍历 {filename}", leave=False):
            page = doc.load_page(page_num)
            
            # 1. 处理文本
            page_text = page.get_text("text")
            if page_text:
                chunks = text_splitter.split_text(page_text)
                for chunk in chunks:
                    text_docs.append(Document(
                        page_content=chunk,
                        metadata={"doc_name": filename, "page": page_num + 1, "clause_id": ""}
                    ))
            
            # 2. 处理图片
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    # (为防止重名, 替换路径中的 / 和 \)
                    safe_filename = filename.replace("\\", "_").replace("/", "_")
                    img_filename = f"{safe_filename}_p{page_num+1}_img{img_index}.{image_ext}"
                    img_save_path = os.path.join(output_image_dir, img_filename)
                    
                    with open(img_save_path, "wb") as f: f.write(image_bytes)
                    
                    image_infos.append({
                        "image_path": img_save_path,
                        "doc_name": filename,
                        "page": page_num + 1
                    })
                except Exception as e:
                    print(f"    > 警告: 提取图片 {xref} 失败: {e}")
                    
        doc.close()
    except Exception as e:
        print(f"  > 错误: 打开 {filename} 失败: {e}")
        
    return text_docs, image_infos

# --- (重大升级) 主函数 ---
def process_all_files() -> Tuple[List[Document], List[Dict]]:
    """
    (主函数 V4)
    递归加载 data/ 目录下的所有文件 (PDF, DOCX, PPTX)，
    并返回文本块和图片信息。
    """
    data_dir = config.DATA_DIR
    output_image_dir = config.OUTPUT_IMAGE_PATH
    os.makedirs(output_image_dir, exist_ok=True)
    
    print(f"--- (V4 递归版) 开始从 {data_dir} 加载所有文件 ---")
    
    all_text_docs: List[Document] = []
    all_image_infos: List[Dict] = []
    
    if not os.path.exists(data_dir):
        print(f"错误: 知识库目录 {data_dir} 不存在。")
        return [], []

    # --- (关键升级) 使用 os.walk 递归遍历 ---
    for root, dirs, files in os.walk(data_dir):
        for filename in files:
            file_path = os.path.join(root, filename)
            
            # (新增) 相对路径名, 用于元数据
            # e.g., "地下室\建筑设计防火规范.pdf"
            relative_filename = os.path.relpath(file_path, data_dir) 

            if filename.endswith(".pdf"):
                text_docs, image_infos = process_pdf(file_path, relative_filename, output_image_dir)
                all_text_docs.extend(text_docs)
                all_image_infos.extend(image_infos)
                
            elif filename.endswith(".docx"):
                text_docs = process_docx(file_path, relative_filename)
                all_text_docs.extend(text_docs)
                
            elif filename.endswith(".pptx"):
                text_docs = process_pptx(file_path, relative_filename)
                all_text_docs.extend(text_docs)
                
            else:
                print(f"  > (跳过) 不支持的文件类型: {filename}")
    
    print(f"--- 所有文件处理完毕 ---")
    print(f"  > 共提取 {len(all_text_docs)} 个文本块")
    print(f"  > 共提取 {len(all_image_infos)} 张图片")
    
    return all_text_docs, all_image_infos